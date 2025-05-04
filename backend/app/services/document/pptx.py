import os
import json
import logging
import mimetypes
from pathlib import Path
from typing import Dict, Any, List, Optional
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ...core.config import settings
from .base import DocumentService
from ..storage import get_storage_service

logger = logging.getLogger(__name__)

class PptxDocumentService(DocumentService):
    """Service for generating PowerPoint presentations."""
    
    def __init__(self):
        self.templates_path = Path(settings.PPTX_TEMPLATES_PATH)
        self.storage_service = get_storage_service()
        
        # Ensure the templates directory exists
        self._ensure_templates_path_exists()
    
    def _ensure_templates_path_exists(self):
        """Ensure that the templates directory exists."""
        if not self.templates_path.exists():
            logger.info(f"Creating templates directory: {self.templates_path}")
            self.templates_path.mkdir(parents=True, exist_ok=True)
    
    def _get_template_path(self, template_name: str) -> Path:
        """Get the full path to a template."""
        template_path = self.templates_path / f"{template_name}.pptx"
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        return template_path
    
    def _get_schema_path(self, template_name: str) -> Path:
        """Get the path to a template's schema file."""
        schema_path = self.templates_path / f"{template_name}.schema.json"
        return schema_path
    
    def _get_info_path(self, template_name: str) -> Path:
        """Get the path to a template's info file."""
        info_path = self.templates_path / f"{template_name}.info.json"
        return info_path
    
    async def generate_document(
        self, 
        template_name: str, 
        data: Dict[str, Any], 
        output_path: Optional[str] = None
    ) -> Path:
        """
        Generate a PowerPoint presentation from a template and data.
        
        Args:
            template_name: Name of the template to use
            data: Data to populate the template with
            output_path: Optional path where to save the document
            
        Returns:
            Path to the generated presentation
        """
        template_path = self._get_template_path(template_name)
        
        # Create a presentation from the template
        prs = Presentation(template_path)
        
        # Process each slide in the presentation
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    self._process_text_frame(shape.text_frame, data)
                if hasattr(shape, "table") and shape.table:
                    self._process_table(shape.table, data)
        
        # Create output filename if not provided
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
            output_path = f"presentations/{template_name}_{timestamp}.pptx"
        
        # Ensure directories exist
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Save the presentation
        prs.save(output_file)
        logger.info(f"Generated presentation: {output_file}")
        
        return output_file
    
    def _process_text_frame(self, text_frame, data: Dict[str, Any]):
        """
        Process a text frame in a slide, replacing placeholders with data.
        
        Args:
            text_frame: The text frame to process
            data: The data dictionary
        """
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                # Replace placeholders of the form {{key}} with data values
                for key, value in data.items():
                    placeholder = f"{{{{{key}}}}}"
                    if placeholder in text:
                        text = text.replace(placeholder, str(value))
                run.text = text
    
    def _process_table(self, table, data: Dict[str, Any]):
        """
        Process a table in a slide, replacing placeholders with data.
        
        Args:
            table: The table to process
            data: The data dictionary
        """
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    self._process_text_frame(cell.text_frame, data)
    
    async def list_templates(self) -> List[Dict[str, Any]]:
        """
        List available PowerPoint templates.
        
        Returns:
            List of template metadata
        """
        templates = []
        
        if not self.templates_path.exists():
            return templates
        
        for file_path in self.templates_path.glob("*.pptx"):
            if file_path.is_file() and not file_path.name.startswith("."):
                template_name = file_path.stem
                
                # Try to get info from the info file
                info = await self.get_template_info(template_name)
                
                templates.append({
                    "name": template_name,
                    "file": file_path.name,
                    "type": "pptx",
                    "info": info
                })
        
        return templates
    
    async def get_template_info(self, template_name: str) -> Dict[str, Any]:
        """
        Get information about a specific template.
        
        Args:
            template_name: Name of the template
            
        Returns:
            Template metadata
        """
        info_path = self._get_info_path(template_name)
        
        # Default info
        info = {
            "name": template_name,
            "type": "pptx",
            "description": f"PowerPoint template: {template_name}",
            "slides": 0
        }
        
        # Try to get info from the info file
        if info_path.exists():
            try:
                with open(info_path, "r") as f:
                    file_info = json.load(f)
                    info.update(file_info)
            except Exception as e:
                logger.error(f"Error reading template info file: {e}")
        
        # Try to get slide count from the template
        try:
            template_path = self._get_template_path(template_name)
            prs = Presentation(template_path)
            info["slides"] = len(prs.slides)
        except Exception as e:
            logger.error(f"Error counting slides in template: {e}")
        
        return info
    
    async def get_template_schema(self, template_name: str) -> Dict[str, Any]:
        """
        Get the schema of data required for a specific template.
        
        Args:
            template_name: Name of the template
            
        Returns:
            JSON schema of the data required for the template
        """
        schema_path = self._get_schema_path(template_name)
        
        # Default schema
        schema = {
            "type": "object",
            "properties": {},
            "required": []
        }
        
        # Try to get schema from the schema file
        if schema_path.exists():
            try:
                with open(schema_path, "r") as f:
                    file_schema = json.load(f)
                    schema.update(file_schema)
            except Exception as e:
                logger.error(f"Error reading template schema file: {e}")
        else:
            # If no schema file exists, try to extract placeholders from the template
            try:
                template_path = self._get_template_path(template_name)
                prs = Presentation(template_path)
                
                # Find all placeholders in the presentation
                placeholders = set()
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text = run.text
                                    # Find all {{key}} patterns
                                    import re
                                    for match in re.finditer(r"{{([^}]+)}}", text):
                                        placeholders.add(match.group(1))
                
                # Create schema properties for each placeholder
                for placeholder in placeholders:
                    schema["properties"][placeholder] = {"type": "string"}
                    schema["required"].append(placeholder)
            
            except Exception as e:
                logger.error(f"Error extracting placeholders from template: {e}")
        
        return schema 