from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
import os
from pathlib import Path
import math

def create_nexus_template():
    """Create an enhanced PPTX template for the Nexus business concept with modern design."""
    
    # Create a blank presentation
    prs = Presentation()
    
    # Define theme colors
    PRIMARY_COLOR = RGBColor(44, 62, 80)     # Dark blue/slate
    SECONDARY_COLOR = RGBColor(52, 152, 219) # Bright blue
    ACCENT_COLOR = RGBColor(231, 76, 60)     # Red/orange
    LIGHT_COLOR = RGBColor(236, 240, 241)    # Light gray
    DARK_COLOR = RGBColor(44, 62, 80)        # Dark blue/slate
    
    # Create a slide master with custom styles
    # Unfortunately, python-pptx doesn't fully support slide master customization
    # So we'll create slides with consistent styling
    
    # Slide 1: Title Slide with Hexagon Theme
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Set background gradient by adding a rectangle
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(41, 128, 185)  # Lighter blue
    fill.gradient_angle = 45  # 45-degree angle
    background.shadow.inherit = False
    background.zorder = 0  # Send to back
    
    # Add hexagon pattern as a visual element
    for i in range(5):
        for j in range(4):
            x = Inches(j * 2 - 1)
            y = Inches(i * 2 - 1)
            hex_shape = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON, x, y, Inches(0.8), Inches(0.8)
            )
            hex_fill = hex_shape.fill
            hex_fill.solid()
            hex_fill.fore_color.rgb = RGBColor(255, 255, 255)
            hex_shape.line.color.rgb = RGBColor(255, 255, 255)
            hex_shape.line.width = Pt(1)
            hex_shape.transparency = 0.9
    
    # Create a central hexagon as the "Nexus" symbol
    central_hex = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON, Inches(4.5-1), Inches(3.5-1), Inches(2), Inches(2)
    )
    central_hex_fill = central_hex.fill
    central_hex_fill.solid()
    central_hex_fill.fore_color.rgb = SECONDARY_COLOR
    central_hex.line.color.rgb = RGBColor(255, 255, 255)
    central_hex.line.width = Pt(2)
    
    # Add data flow lines converging into the central hexagon
    data_flow_points = [
        (Inches(1), Inches(1)),       # Top left
        (Inches(1), Inches(6)),       # Bottom left
        (Inches(5), Inches(0.5)),     # Top middle
        (Inches(9), Inches(1)),       # Top right
        (Inches(9), Inches(6)),       # Bottom right
        (Inches(5), Inches(6.5)),     # Bottom middle
    ]
    
    # Central point of the hexagon
    center_x = Inches(4.5-1) + Inches(1)  # Left + half width
    center_y = Inches(3.5-1) + Inches(1)  # Top + half height
    
    # Add converging lines and small data point circles
    for x, y in data_flow_points:
        # Add a small circle representing a data point
        data_point = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4)
        )
        point_fill = data_point.fill
        point_fill.solid()
        
        # Use different colors for variety
        if x < Inches(5):  # Left side
            point_fill.fore_color.rgb = ACCENT_COLOR
        else:  # Right side
            point_fill.fore_color.rgb = RGBColor(46, 204, 113)  # Green
            
        data_point.line.color.rgb = RGBColor(255, 255, 255)
        data_point.line.width = Pt(1)
        
        # Add "data flowing" connector from point to center
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x + Inches(0.2),  # Center of data point
            y + Inches(0.2),
            center_x,  # Center of hexagon
            center_y
        )
        connector.line.color.rgb = RGBColor(255, 255, 255)
        connector.line.width = Pt(1.5)
        
        # Add small arrows along the connector to show flow direction
        # (We'll simulate this with small triangle shapes)
        # Calculate midpoint of connector
        mid_x = (x + Inches(0.2) + center_x) / 2
        mid_y = (y + Inches(0.2) + center_y) / 2
        
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE, mid_x - Inches(0.15), mid_y - Inches(0.15), Inches(0.3), Inches(0.3)
        )
        arrow_fill = arrow.fill
        arrow_fill.solid()
        arrow_fill.fore_color.rgb = RGBColor(255, 255, 255)
        arrow.rotation = 180  # Point toward the center
        
        # Adjust rotation to point toward center
        dx = center_x - (x + Inches(0.2))
        dy = center_y - (y + Inches(0.2))
        angle = math.degrees(math.atan2(dy, dx))
        arrow.rotation = angle
    
    # Style the title and subtitle
    title.text = "Nexus Solutions"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    title.top = Inches(1)
    title.left = Inches(5)
    title.width = Inches(4)
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.size = Pt(44)
    
    subtitle.text = "Where Data Converges, Insight Emerges\n\n{{current_date}}"
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    subtitle.top = Inches(2.2)
    subtitle.left = Inches(5)
    subtitle.width = Inches(4)
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = LIGHT_COLOR
            run.font.size = Pt(20)
    
    # Slide 2: Nexus Defined with Visual Elements
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    # Add a subtle background element
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 0, Inches(0.5), prs.slide_width, Inches(0.8)
    )
    bg_fill = bg.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = SECONDARY_COLOR
    bg.line.fill.background()  # No line
    bg.shadow.inherit = False
    bg.zorder = 0  # Send to back
    
    # Style the title
    title.text = "Nexus Defined"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.size = Pt(36)
    
    # Add a visual dictionary definition
    dict_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(4), Inches(3)
    )
    dict_fill = dict_box.fill
    dict_fill.solid()
    dict_fill.fore_color.rgb = LIGHT_COLOR
    dict_box.line.color.rgb = SECONDARY_COLOR
    dict_box.line.width = Pt(2)
    dict_box.shadow.inherit = False
    
    # Add text to the definition box
    text_frame = dict_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "Nexus (noun):"
    p = text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    
    p = text_frame.add_paragraph()
    p.text = "1. A connection or series of connections linking two or more things."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_COLOR
    
    p = text_frame.add_paragraph()
    p.text = "2. A central or focal point."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_COLOR
    
    p = text_frame.add_paragraph()
    p.text = ""
    
    p = text_frame.add_paragraph()
    p.text = "At Nexus Solutions, we are the central connection point for all your data needs."
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Add a visual representation of "nexus" concept
    for i, color in enumerate([
        SECONDARY_COLOR,
        PRIMARY_COLOR,
        ACCENT_COLOR,
        RGBColor(46, 204, 113),  # Green
        RGBColor(155, 89, 182)   # Purple
    ]):
        angle = i * 72  # Distribute 5 circles around center
        x = Inches(7) + Inches(1.2) * (angle % 180) / 180
        y = Inches(3.5) + Inches(0.8) * ((angle + 90) % 180) / 180
        
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6)
        )
        circle_fill = circle.fill
        circle_fill.solid()
        circle_fill.fore_color.rgb = color
        
        # Add connection line to center
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, x + Inches(0.3), y + Inches(0.3), 
            Inches(8), Inches(3.5)
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(2)
        
    # Add central circle
    center_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.7), Inches(3.2), Inches(0.6), Inches(0.6)
    )
    center_fill = center_circle.fill
    center_fill.solid()
    center_fill.fore_color.rgb = PRIMARY_COLOR
    center_circle.line.color.rgb = RGBColor(255, 255, 255)
    center_circle.line.width = Pt(2)
    
    text_frame = center_circle.text_frame
    text_frame.text = "N"
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide 3: Business Overview with Visual Hierarchy
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    
    # Add a subtle background element
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 0, Inches(0.5), prs.slide_width, Inches(0.8)
    )
    bg_fill = bg.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = SECONDARY_COLOR
    bg.line.fill.background()  # No line
    bg.shadow.inherit = False
    bg.zorder = 0  # Send to back
    
    # Style the title
    title.text = "Business Overview"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.size = Pt(36)
    
    # Add custom content with better visual hierarchy
    overview_text = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1)
    )
    text_frame = overview_text.text_frame
    text_frame.word_wrap = True
    text_frame.text = "{{business_name}} provides advanced data solutions for businesses and individuals with complex data needs."
    p = text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_COLOR
    
    # Create a more visual list of features using shapes
    features = [
        "Industry-specific vertical solutions",
        "Customized integration of diverse data sources",
        "Powerful visualization and analytics",
        "Scalable architecture for businesses of all sizes",
        "{{unique_value_proposition}}"
    ]
    
    for i, feature in enumerate(features):
        y_pos = Inches(3 + i * 0.7)
        
        # Add feature icon
        icon = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y_pos, Inches(0.4), Inches(0.4)
        )
        icon_fill = icon.fill
        icon_fill.solid()
        icon_fill.fore_color.rgb = SECONDARY_COLOR
        icon.line.fill.background()  # No line
        
        # Add feature text
        feature_text = slide.shapes.add_textbox(
            Inches(1.6), y_pos, Inches(7), Inches(0.6)
        )
        text_frame = feature_text.text_frame
        text_frame.word_wrap = True
        text_frame.text = feature
        p = text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_COLOR
    
    # Slide 4: Value Proposition with Data Convergence Theme
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Title and Two Content
    title = slide.shapes.title
    
    # Add a subtle background element
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 0, Inches(0.5), prs.slide_width, Inches(0.8)
    )
    bg_fill = bg.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = SECONDARY_COLOR
    bg.line.fill.background()  # No line
    bg.shadow.inherit = False
    bg.zorder = 0  # Send to back
    
    # Style the title
    title.text = "Our Value Proposition"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.size = Pt(36)
    
    # Left: Value propositions in a convergence flow diagram
    values = [
        "Single source of truth for all data",
        "Streamlined data workflows",
        "Reduced complexity",
        "Real-time analytics",
        "Data-driven decision making"
    ]
    
    # Create a convergence flow diagram with better spacing
    # Start with a large "Insights" circle at the bottom
    insights_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(2.5), Inches(4.5), Inches(1.5), Inches(1.5)
    )
    insights_fill = insights_circle.fill
    insights_fill.solid()
    insights_fill.fore_color.rgb = SECONDARY_COLOR
    insights_circle.line.color.rgb = RGBColor(255, 255, 255)
    insights_circle.line.width = Pt(2)
    
    # Add text to the insights circle
    insights_text = insights_circle.text_frame
    insights_text.text = "INSIGHTS"
    insights_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    insights_text.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = insights_text.paragraphs[0].runs[0]
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add value proposition elements in a fan above the insights circle
    fan_radius = Inches(3)
    center_x = Inches(3.25)  # Center of the insights circle
    center_y = Inches(5.25)  # Center of the insights circle
    
    # Calculate positions in an arc above the insights circle
    for i, value in enumerate(values):
        # Calculate position in an arc
        angle = -60 + (i * 30)  # Spread from -60 to +60 degrees
        angle_rad = math.radians(angle)
        
        # Position above the insights circle at varying distances
        x = center_x + fan_radius * math.sin(angle_rad) - Inches(1)
        y = center_y - fan_radius * math.cos(angle_rad) - Inches(0.5)
        
        # Add value proposition in a rounded rectangle
        value_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2), Inches(0.75)
        )
        value_fill = value_box.fill
        value_fill.solid()
        value_fill.fore_color.rgb = LIGHT_COLOR
        value_box.line.color.rgb = PRIMARY_COLOR
        value_box.line.width = Pt(1.5)
        
        # Add text to the value box
        text_frame = value_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = value
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = DARK_COLOR
        
        # Add flow line from value to insights
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x + Inches(1),  # Middle of value box
            y + Inches(0.75),  # Bottom of value box
            center_x,  # Top of insights circle
            center_y - Inches(0.5)
        )
        connector.line.color.rgb = PRIMARY_COLOR
        connector.line.width = Pt(2)
    
    # Right: Visual representation of Convergence
    # Create a central "Nexus" circle
    central_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7), Inches(3.5), Inches(1.2), Inches(1.2)
    )
    circle_fill = central_circle.fill
    circle_fill.solid()
    circle_fill.fore_color.rgb = SECONDARY_COLOR
    central_circle.line.color.rgb = RGBColor(255, 255, 255)
    central_circle.line.width = Pt(2)
    
    # Add "NEXUS" text to circle
    text_frame = central_circle.text_frame
    text_frame.text = "NEXUS"
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add data sources flowing into the central circle
    data_sources = [
        ("CRM Data", Inches(6), Inches(1.5), ACCENT_COLOR),
        ("ERP Data", Inches(8.5), Inches(1.5), RGBColor(46, 204, 113)),  # Green
        ("IoT Data", Inches(9.5), Inches(3.5), RGBColor(155, 89, 182)),  # Purple
        ("Web Data", Inches(8.5), Inches(5.5), RGBColor(243, 156, 18)),  # Orange
        ("API Data", Inches(6), Inches(5.5), RGBColor(52, 152, 219)),    # Blue
        ("DB Data", Inches(5), Inches(3.5), RGBColor(231, 76, 60))       # Red
    ]
    
    for label, x, y, color in data_sources:
        # Add source rectangle
        source = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1), Inches(0.5)
        )
        source_fill = source.fill
        source_fill.solid()
        source_fill.fore_color.rgb = color
        source.line.color.rgb = RGBColor(255, 255, 255)
        source.line.width = Pt(1)
        
        # Add text to source
        text_frame = source.text_frame
        text_frame.text = label
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add flow connector from source to Nexus
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x + Inches(0.5),  # Center of source
            y + Inches(0.25),
            Inches(7.6),  # Center of Nexus circle
            Inches(4.1)
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(1.5)
        
        # Add small data packet symbols along the connector
        # Calculate midpoint of connector
        mid_x = (x + Inches(0.5) + Inches(7.6)) / 2
        mid_y = (y + Inches(0.25) + Inches(4.1)) / 2
        
        packet = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, mid_x - Inches(0.1), mid_y - Inches(0.1), Inches(0.2), Inches(0.2)
        )
        packet_fill = packet.fill
        packet_fill.solid()
        packet_fill.fore_color.rgb = color
    
    # Additional slides would follow similar pattern with consistent styling
    # For brevity, we'll add one more slide with financial data
    
    # Slide 5: Financial Projections with growth curve visualization
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    
    # Add a subtle background element
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 0, Inches(0.5), prs.slide_width, Inches(0.8)
    )
    bg_fill = bg.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = SECONDARY_COLOR
    bg.line.fill.background()  # No line
    bg.shadow.inherit = False
    bg.zorder = 0  # Send to back
    
    # Style the title
    title.text = "Financial Projections"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.size = Pt(36)
    
    # Add subtitle text
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(0.5)
    )
    text_frame = subtitle_box.text_frame
    text_frame.text = "Where Financial Data Converges into Strategic Growth"
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Create a cleaner, more professional table for projections
    # Add table header with better spacing
    headers = ["Metric", "Year 1", "Year 3", "Year 5"]
    header_y = Inches(2.2)
    
    # Create table background
    table_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(8), Inches(2.5)
    )
    table_bg_fill = table_bg.fill
    table_bg_fill.solid()
    table_bg_fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray
    table_bg.line.color.rgb = RGBColor(230, 230, 230)
    table_bg.line.width = Pt(1)
    table_bg.shadow.inherit = False
    
    # Position headers with better spacing
    for i, header in enumerate(headers):
        x = Inches(1.5 + i * 2)
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, header_y, Inches(1.8), Inches(0.4)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = PRIMARY_COLOR
        header_shape.line.fill.background()
        header_shape.shadow.inherit = False
        
        text_frame = header_shape.text_frame
        text_frame.text = header
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add table rows with consistent spacing
    metrics = ["Revenue", "Customers", "Team Size"]
    placeholder_values = [
        ["{{year1_revenue}}", "{{year3_revenue}}", "{{year5_revenue}}"],
        ["{{year1_customers}}", "{{year3_customers}}", "{{year5_customers}}"],
        ["{{year1_team_size}}", "{{year3_team_size}}", "{{year5_team_size}}"]
    ]
    
    # Use better spacing for rows
    for row, (metric, values) in enumerate(zip(metrics, placeholder_values)):
        row_y = Inches(2.7 + row * 0.6)
        
        # Add metric label with consistent sizing
        metric_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.5), row_y, Inches(1.8), Inches(0.4)
        )
        metric_fill = metric_shape.fill
        metric_fill.solid()
        metric_fill.fore_color.rgb = LIGHT_COLOR
        metric_shape.line.color.rgb = RGBColor(210, 210, 210)
        metric_shape.line.width = Pt(1)
        metric_shape.shadow.inherit = False
        
        text_frame = metric_shape.text_frame
        text_frame.text = metric
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = DARK_COLOR
        
        # Add values with consistent sizing and alignment
        for i, value in enumerate(values):
            x = Inches(3.5 + i * 2)
            value_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, row_y, Inches(1.8), Inches(0.4)
            )
            value_fill = value_shape.fill
            value_fill.solid()
            value_fill.fore_color.rgb = LIGHT_COLOR
            value_shape.line.color.rgb = RGBColor(210, 210, 210)
            value_shape.line.width = Pt(1)
            value_shape.shadow.inherit = False
            
            text_frame = value_shape.text_frame
            text_frame.text = value
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            run = text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.color.rgb = DARK_COLOR
    
    # Add growth curve visualization
    # Create axis lines for the growth chart
    chart_x = Inches(1.5)
    chart_y = Inches(4.8)
    chart_width = Inches(7.5)
    chart_height = Inches(1.5)
    
    # X-axis (horizontal line)
    x_axis = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        chart_x, chart_y + chart_height,  # Start at bottom left
        chart_x + chart_width, chart_y + chart_height  # End at bottom right
    )
    x_axis.line.color.rgb = DARK_COLOR
    x_axis.line.width = Pt(1.5)
    
    # Y-axis (vertical line)
    y_axis = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        chart_x, chart_y + chart_height,  # Start at bottom left
        chart_x, chart_y  # End at top left
    )
    y_axis.line.color.rgb = DARK_COLOR
    y_axis.line.width = Pt(1.5)
    
    # Add year labels on x-axis
    years = ["Year 1", "Year 3", "Year 5"]
    for i, year in enumerate(years):
        # Position evenly along x-axis
        x = chart_x + (i * chart_width / 2)
        y = chart_y + chart_height + Inches(0.1)
        
        year_label = slide.shapes.add_textbox(
            x, y, Inches(1), Inches(0.3)
        )
        text_frame = year_label.text_frame
        text_frame.text = year
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_COLOR
    
    # Create growth curves with different colors for each metric
    curve_data = [
        ("Revenue", SECONDARY_COLOR, [(0, 0.2), (0.5, 0.5), (1, 0.9)]),  # Blue curve
        ("Customers", ACCENT_COLOR, [(0, 0.15), (0.5, 0.45), (1, 0.85)]),  # Red curve
        ("Team", RGBColor(46, 204, 113), [(0, 0.1), (0.5, 0.4), (1, 0.8)])  # Green curve
    ]
    
    for name, color, points in curve_data:
        # Convert relative coordinates to actual positions
        curve_points = []
        for rx, ry in points:
            x = chart_x + (rx * chart_width)
            # Invert y coordinate (0 is bottom, 1 is top)
            y = chart_y + chart_height - (ry * chart_height)
            curve_points.append((x, y))
        
        # Draw curve segments
        for i in range(len(curve_points) - 1):
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                curve_points[i][0], curve_points[i][1],
                curve_points[i+1][0], curve_points[i+1][1]
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(3)
        
        # Add points along the curve
        for x, y in curve_points:
            point = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x - Inches(0.1), y - Inches(0.1), Inches(0.2), Inches(0.2)
            )
            point_fill = point.fill
            point_fill.solid()
            point_fill.fore_color.rgb = color
            point.line.color.rgb = RGBColor(255, 255, 255)
            point.line.width = Pt(1)
        
        # Add curve label at the end of the curve
        label = slide.shapes.add_textbox(
            curve_points[-1][0] + Inches(0.1),
            curve_points[-1][1] - Inches(0.15),
            Inches(1), Inches(0.3)
        )
        text_frame = label.text_frame
        text_frame.text = name
        run = text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = color
    
    # Add investment information box at the bottom
    investment_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.5), Inches(5.5), Inches(0.8)
    )
    investment_fill = investment_box.fill
    investment_fill.solid()
    investment_fill.fore_color.rgb = PRIMARY_COLOR
    investment_box.line.fill.background()
    investment_box.shadow.inherit = False
    
    text_frame = investment_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = "Initial Investment: "
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    run = p.add_run()
    run.text = "{{initial_investment}}  •  "
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    run = p.add_run()
    run.text = "Breakeven Timeline: "
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    run = p.add_run()
    run.text = "{{breakeven_timeline}}"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    # Create additional template slides following similar pattern
    # For brevity, we'll add just one more - the final "Thank You" slide
    
    # Slide 11: Thank You with Data Convergence Theme
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Set background gradient by adding a rectangle
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[0].color.rgb = PRIMARY_COLOR
    fill.gradient_stops[1].position = 1
    fill.gradient_stops[1].color.rgb = RGBColor(41, 128, 185)  # Lighter blue
    fill.gradient_angle = 45  # 45-degree angle
    background.shadow.inherit = False
    background.zorder = 0  # Send to back
    
    # Create a visually appealing data convergence illustration
    # Add flowing data streams that converge in the center
    
    # Center point where all data converges
    center_x = Inches(5)
    center_y = Inches(3.5)
    
    # Radial lines representing data streams
    num_streams = 12
    stream_length = Inches(4)
    
    for i in range(num_streams):
        angle = i * (360 / num_streams)
        angle_rad = math.radians(angle)
        
        # Calculate outer point of the stream
        outer_x = center_x + stream_length * math.cos(angle_rad)
        outer_y = center_y + stream_length * math.sin(angle_rad)
        
        # Draw data stream
        stream = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            outer_x, outer_y,
            center_x, center_y
        )
        
        # Alternate colors for visual interest
        if i % 3 == 0:
            stream_color = SECONDARY_COLOR  # Blue
        elif i % 3 == 1:
            stream_color = ACCENT_COLOR  # Red
        else:
            stream_color = RGBColor(46, 204, 113)  # Green
            
        stream.line.color.rgb = stream_color
        stream.line.width = Pt(2)
        
        # Add data packet symbols along the stream
        for j in range(1, 4):  # 3 packets per stream
            # Position along the stream
            packet_x = outer_x - (j * (stream_length / 4) * math.cos(angle_rad))
            packet_y = outer_y - (j * (stream_length / 4) * math.sin(angle_rad))
            
            # Alternate packet shapes
            if j % 2 == 0:
                packet = slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND, packet_x - Inches(0.1), packet_y - Inches(0.1), 
                    Inches(0.2), Inches(0.2)
                )
            else:
                packet = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, packet_x - Inches(0.1), packet_y - Inches(0.1), 
                    Inches(0.2), Inches(0.2)
                )
                
            packet_fill = packet.fill
            packet_fill.solid()
            packet_fill.fore_color.rgb = stream_color
            packet.line.color.rgb = RGBColor(255, 255, 255)
            packet.line.width = Pt(1)
    
    # Create a central "insight" glow effect
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, center_x - Inches(1.2), center_y - Inches(1.2), 
        Inches(2.4), Inches(2.4)
    )
    glow_fill = glow.fill
    glow_fill.solid()
    glow_fill.fore_color.rgb = RGBColor(255, 255, 255)
    glow.line.fill.background()
    
    # Add the thank you text overlaid on the glow
    thank_you = slide.shapes.add_textbox(
        center_x - Inches(2.5), center_y - Inches(1.5), 
        Inches(5), Inches(1)
    )
    thank_text = thank_you.text_frame
    thank_text.text = "THANK YOU"
    thank_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = thank_text.paragraphs[0].runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_COLOR
    
    # Add business name and tagline
    tagline_box = slide.shapes.add_textbox(
        center_x - Inches(3.5), center_y - Inches(0.25), 
        Inches(7), Inches(1)
    )
    tagline_text = tagline_box.text_frame
    
    p = tagline_text.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = "{{business_name}}"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    
    p = tagline_text.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(24)
    
    run = p.add_run()
    run.text = "{{tagline}}"
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = LIGHT_COLOR
    
    # Add contact information at bottom
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        center_x - Inches(2.5), center_y + Inches(1.5), 
        Inches(5), Inches(1.5)
    )
    contact_fill = contact_box.fill
    contact_fill.solid()
    contact_fill.fore_color.rgb = RGBColor(255, 255, 255)
    contact_fill.transparency = 0.2  # Semi-transparent
    contact_box.line.color.rgb = RGBColor(255, 255, 255)
    contact_box.line.width = Pt(1)
    
    contact_text = contact_box.text_frame
    
    p = contact_text.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(6)
    
    run = p.add_run()
    run.text = "Contact:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_COLOR
    
    p = contact_text.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    
    run = p.add_run()
    run.text = "{{contact_name}}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK_COLOR
    
    p = contact_text.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    
    run = p.add_run()
    run.text = "{{contact_email}}"
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_COLOR
    
    p = contact_text.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = "{{contact_phone}}"
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_COLOR
    
    # Create output directory if it doesn't exist
    template_dir = Path("data/templates/pptx")
    template_dir.mkdir(parents=True, exist_ok=True)
    
    # Save the presentation
    output_path = template_dir / "nexus_concept_enhanced.pptx"
    prs.save(output_path)
    print(f"Enhanced template created successfully: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_nexus_template() 